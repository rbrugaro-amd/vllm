"""Generate the vLLM torch.compile & Fusions deck on the AMD corporate template.

Reproducible: reads `corporate_template.pptx` (converted from the .potx), strips the
bundled demo slides, and builds the deck from the template's own layouts/master so
all AMD branding, fonts, and colors are inherited.

Run:  .slides-venv/bin/python gen_slides.py
Out:  vllm_compile_fusions_amd.pptx
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

TEMPLATE = "corporate_template.pptx"
OUTPUT = "vllm_compile_fusions_amd.pptx"

# AMD palette (from template theme)
AMD_ORANGE = RGBColor(0xF2, 0x65, 0x22)
AMD_CYAN = RGBColor(0x00, 0xC2, 0xDE)
AMD_GRAY = RGBColor(0x5E, 0x5E, 0x5E)
MONO = "Consolas"

# Layout indices (from inspect_template.py)
L_TITLE = 0          # Title Slide - No Image
L_TITLE_CONTENT = 3  # Title and Content
L_TWO_CONTENT = 5    # Two Content
L_SECTION = 8        # Small title - half page (section divider)
L_AGENDA = 22        # Agenda
L_CODE = 27          # Developer Code Layout
L_SCORECARD = 20     # Scorecard
L_TITLE_ONLY = 7     # Title Only


def delete_demo_slides(prs):
    """Remove bundled demo slides AND drop their relationships so the orphaned
    slide/media parts are not serialized (avoids duplicate partnames, slims file)."""
    part = prs.part
    lst = prs.slides._sldIdLst
    for sldId in list(lst):
        rId = sldId.get(qn("r:id"))
        part.drop_rel(rId)
        lst.remove(sldId)


def ph(slide, idx):
    for p in slide.placeholders:
        if p.placeholder_format.idx == idx:
            return p
    raise KeyError(f"placeholder idx={idx} not found in layout "
                   f"'{slide.slide_layout.name}'")


def title_ph(slide):
    return slide.shapes.title


def _emit(tf, items, clear=True):
    """items: list of (text, level) or (text, level, opts) where opts is a dict.
    opts keys: bold(bool), color(RGBColor), mono(bool), size(pt int)."""
    if clear:
        tf.clear()
    first = True
    for item in items:
        text, level = item[0], item[1]
        opts = item[2] if len(item) > 2 else {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        runs = text if isinstance(text, list) else [(text, {})]
        for rtext, ropts in runs:
            r = p.add_run()
            r.text = rtext
            merged = {**opts, **ropts}
            if merged.get("bold"):
                r.font.bold = True
            if merged.get("color") is not None:
                r.font.color.rgb = merged["color"]
            if merged.get("mono"):
                r.font.name = MONO
            if merged.get("size"):
                r.font.size = Pt(merged["size"])


def add_section(prs, title):
    s = prs.slides.add_slide(prs.slide_layouts[L_SECTION])
    title_ph(s).text = title
    return s


def add_bullets(prs, title, items):
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_CONTENT])
    title_ph(s).text = title
    _emit(ph(s, 10).text_frame, items)
    return s


def add_takeaway(slide, items):
    """Add a full-width takeaway band near the bottom of a slide."""
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    _emit(tf, items)
    for p in tf.paragraphs:
        for r in p.runs:
            if not r.font.size:
                r.font.size = Pt(13)
    return box


def _no_bullet(p):
    pPr = p._pPr
    if pPr is None:
        pPr = p.get_or_add_pPr()
    # remove existing bullet defs, then add explicit none
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def add_code(prs, title, code, note=None, size=11):
    s = prs.slides.add_slide(prs.slide_layouts[L_CODE])
    title_ph(s).text = title
    tf = ph(s, 10).text_frame
    tf.word_wrap = True
    tf.clear()
    lines = code.strip("\n").split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        _no_bullet(p)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = MONO
        r.font.size = Pt(size)
        if line.lstrip().startswith("#"):
            r.font.color.rgb = AMD_GRAY
    if note:
        add_takeaway(s, [([o(note)], 0)])
    return s


def add_image_slide(prs, title, img_path, caption=None, top=1.9, width_in=12.4):
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
    title_ph(s).text = title
    left = Inches((13.33 - width_in) / 2)
    s.shapes.add_picture(img_path, left, Inches(top), width=Inches(width_in))
    if caption:
        add_takeaway(s, [([o(caption)], 0)])
    return s


def add_table_slide(prs, title, headers, rows, note=None):
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
    title_ph(s).text = title
    nrows, ncols = len(rows) + 1, len(headers)
    gx = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(1.7),
                            Inches(12.1), Inches(0.4 * nrows)).table
    for c, h in enumerate(headers):
        cell = gx.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(13)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = gx.cell(ri, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
    if note:
        add_takeaway(s, [([o(note)], 0)])
    return s


def add_code_explain(prs, title, code, explain, code_size=9):
    """Two-column slide: monospace code on the left, explanation bullets right."""
    s = prs.slides.add_slide(prs.slide_layouts[L_TWO_CONTENT])
    title_ph(s).text = title
    tf = ph(s, 12).text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(code.strip("\n").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        _no_bullet(p)
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = MONO
        r.font.size = Pt(code_size)
        if line.lstrip().startswith("#"):
            r.font.color.rgb = AMD_GRAY
    _emit(ph(s, 13).text_frame, explain)
    return s


def add_two_content(prs, title, left, right):
    s = prs.slides.add_slide(prs.slide_layouts[L_TWO_CONTENT])
    title_ph(s).text = title
    _emit(ph(s, 12).text_frame, left)
    _emit(ph(s, 13).text_frame, right)
    return s


# emphasis helpers for inline runs
def b(t):
    return (t, {"bold": True})


def o(t):
    return (t, {"bold": True, "color": AMD_ORANGE})


def c(t):
    return (t, {"mono": True})


def build():
    prs = Presentation(TEMPLATE)
    delete_demo_slides(prs)

    # ---- Title ----
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE])
    title_ph(s).text = "vLLM torch.compile & Fusions on MI355"
    _emit(ph(s, 13).text_frame, [
        ("Bringing performance uplifts to vLLM through the compiler", 0),
        ("A deep dive, from concepts to a real fusion pass (PR #39242)", 0),
    ])

    # ---- Agenda ----
    s = prs.slides.add_slide(prs.slide_layouts[L_AGENDA])
    title_ph(s).text = "Agenda"
    _emit(ph(s, 10).text_frame, [
        ("Part I — Foundations: CUDA Graphs vs torch.compile, vLLM philosophy", 0),
        ("Part II — How vLLM compilation works: piecewise, the flow, vLLM IR, nuances", 0),
        ("Part III — Worked example: MLA dual RMSNorm fusion (PR #39242)", 0),
        ("Part IV — Practical: compile config & debug BKM", 0),
    ])

    # =====================================================================
    # PART I — FOUNDATIONS
    # =====================================================================
    add_section(prs, "Part I — Foundations")

    # S1: CUDA Graphs vs torch.compile (with merged "complementary" takeaway band)
    s = add_two_content(
        prs,
        "Two problems, two tools",
        left=[
            ([o("CUDA Graphs")], 0),
            ("Attacks CPU launch overhead", 1),
            ("Capture a launch sequence once, replay with ~zero CPU cost", 1),
            ("Constraints: static shapes, fixed addresses, GPU-only (no CPU branching)", 1),
        ],
        right=[
            ([o("torch.compile / Inductor")], 0),
            ("Attacks kernel quality & memory traffic", 1),
            ("Generates fused Triton/C++ kernels, fewer launches", 1),
            ("Can autotune; cuts full-precision memory round-trips", 1),
        ],
    )
    add_takeaway(s, [
        ([b("Complementary, not either/or"), (" — vLLM layers both: Inductor-compiled subgraphs replayed inside CUDA graphs", {})], 0),
        ([o("Thesis: make the model fast by teaching the compiler, not rewriting the model")], 0),
    ])

    # S2: vLLM philosophy (with merged, condensed SGLang contrast)
    add_bullets(prs, "vLLM optimization philosophy", [
        ([b("Separation of concerns")], 0),
        ("Model code stays clean — math only; optimizations live in compiler passes", 1),
        ([b("Why it matters")], 0),
        ("One pass benefits every model that hits the pattern — no per-model fused-kernel copies, abstractions intact", 1),
        ([("Contrast with ", {}), b("SGLang"), (": leans on hand-written / templated fused kernels + model-level integration", {})], 0),
        ("vLLM bets on a compiler pipeline + pattern-matching passes (generality & maintainability vs. per-case control)", 1),
        ([("This is exactly why AMD work ships as a ", {}), o("pass"), (" (PR #39242), not a model edit", {})], 0),
    ])

    # S2c: PassConfig & optimization levels
    add_bullets(prs, "PassConfig & optimization levels", [
        ([b("Levels are presets of underlying flags; user-set flags always win")], 0),
        ("-O0 nothing · -O1 compile + PIECEWISE graphs + fast fusions", 1),
        ("-O2 (default) + FULL_AND_PIECEWISE + more fusions/ranges · -O3 = O2 today", 1),
        ([b("Three buckets of fusions")], 0),
        ([("Always-on at a level", {"bold": True}), (" (e.g. fuse_allreduce_rms at O2 on Hopper/Blackwell)", {})], 1),
        ([("Conditionally auto-enabled", {"bold": True}), (" — platform / custom-op guards", {})], 1),
        ([("e.g. ", {}), c("fuse_mla_dual_rms_norm"), (" only at O1+ on ROCm/AITER", {})], 2),
        ([o("Never auto — opt-in only"), (" — no default at any level", {})], 1),
        ([("e.g. ", {}), c("fuse_attn_quant"), (", ", {}), c("enable_sp"), (", ", {}), c("enable_qk_norm_rope_fusion"), ("  → must force via config", {})], 2),
    ])

    # =====================================================================
    # PART II — HOW vLLM COMPILATION WORKS
    # =====================================================================
    add_section(prs, "Part II — How vLLM Compilation Works")

    # S3: piecewise compile vs piecewise graph
    s = add_two_content(
        prs,
        "Piecewise compilation vs piecewise CUDA graph",
        left=[
            ([b("Unifying idea: split the full graph at attention ops (splitting_ops)")], 0),
            ([o("Piecewise compilation")], 0),
            ("Inductor compiles each inter-attention subgraph independently", 1),
            ("Attention is an opaque custom op", 1),
        ],
        right=[
            ([o("Piecewise CUDA graph")], 0),
            ("Capture graphs only for the compute-between-attentions pieces", 1),
            ("Attention runs eager → stays flexible", 1),
            ([b("Why attention is the seam")], 0),
            ("Output shares the query's shape; between-attn work is token-wise", 1),
        ],
    )
    add_takeaway(s, [
        ([b("Two independent size lists:"), (" ", {}), ("compile_sizes", {"mono": True}), (" → Inductor specialized + autotuned kernels;  ", {}), ("cudagraph_capture_sizes", {"mono": True}), (" → graphs captured, real batches padded up", {})], 0),
    ])

    # S3b: attention is special on two axes (merged with the AITER MLA example)
    add_bullets(prs, "Attention is special: two axes (with AITER MLA)", [
        ([("Axis A — ", {}), o("Inductor compilation"), (": attention is ALWAYS excluded (opaque custom op; Dynamo won't trace in)", {})], 0),
        ([("Axis B — ", {}), o("CUDA graph capture"), (": conditional on backend capability + batch shape", {})], 0),
        ([b("Concrete (resolves \"but I see AITER MLA inside a CUDA graph!\")")], 0),
        ([("Default ", {}), ("FULL_AND_PIECEWISE", {"mono": True}), (": full graph for uniform decode, piecewise (attn eager) for prefill/mixed", {})], 1),
        ([("AITER MLA = ", {}), ("UNIFORM_SINGLE_TOKEN_DECODE", {"mono": True}), (" → captured in full graph during decode; eager for prefill/mixed (both true)", {})], 1),
        ([o("Can't"), (": capture ONE static graph valid across prefill/mixed — fixed addresses + launch geometry; prefill shapes vary per step", {})], 0),
        ([("Ladder: ", {}), ("ALWAYS", {"mono": True}), (" > ", {}), ("UNIFORM_BATCH", {"mono": True}), (" > ", {}), ("UNIFORM_SINGLE_TOKEN_DECODE", {"mono": True}), (" > ", {}), ("NEVER", {"mono": True})], 0),
    ])

    # S4: the compilation flow
    # Combined: official diagram on top + detailed stage mapping beneath
    s = prs.slides.add_slide(prs.slide_layouts[L_TITLE_ONLY])
    title_ph(s).text = "The compilation flow"
    s.shapes.add_picture(
        "docs/assets/design/debug_vllm_compile/design_diagram.png",
        Inches(1.27), Inches(1.3), width=Inches(10.8))
    box = s.shapes.add_textbox(Inches(0.5), Inches(4.95), Inches(12.3), Inches(2.4))
    tf = box.text_frame
    tf.word_wrap = True
    _emit(tf, [
        ([b("Mental model: "), ("Dynamo captures → ", {}), b("vLLM backend splits"), (" → Inductor compiles each piece → vLLM wraps in CUDA graphs", {})], 0),
        ([("1 capture → ", {}), o("Dynamo"), ("    2 splitting → ", {}), o("vLLM backend"), (" (split @ attention)    4 wrapper → ", {}), o("CUDAGraph replay"), (" (attn eager)", {})], 0),
        ([("3 \"backend compile\" (per piece): ", {}),
          ("pre_grad_custom_pass", {"mono": True}), (" → ", {}),
          ("AOTAutograd", {}), (" → ", {}),
          ("post_grad_custom_post_pass", {"mono": True}), (" (", {}),
          b("fusions + IR lowering"), (") → ", {}),
          o("Inductor codegen"), (" (Triton)", {})], 0),
        ([o("backend compile → ONE artifact / piece"), (" = a ", {}), b("Python wrapper"), (" launching Triton ", {}), ("+", {}), (" extern kernels (AITER / CK / cuBLAS); cached, then the CUDAGraph wrapper replays it", {})], 0),
    ])
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(12)

    add_bullets(prs, "Key facts about vLLM's compiler", [
        ([o("vLLM-compile is NOT stock torch.compile")], 0),
        ("Custom compiler on internal PyTorch Compile APIs: full-graph capture,", 1),
        ("custom split, custom passes, own cache", 1),
        ([("Fusions live in ", {}), b("pre-grad and post-grad custom passes"), ("  ← the hook for Part III", {})], 0),
        ([b("Compilation cache")], 0),
        ("Hash factors: configs + PyTorch config + traced files; cache dir is copyable", 1),
        ("Guarantee: all compilation finishes before serving any request", 1),
        ([("Only symbolic dim is ", {}), o("batch size (#tokens)"), ("; weights/buffers are static", {})], 0),
    ])

    # S5: vLLM IR (what/why + why-it-matters, merged)
    add_bullets(prs, "vLLM IR — what it is & why it matters", [
        ([b("The gap it fills:")], 0),
        ("Functional FX 'dialect' between low-level torch ops and vLLM layers (RMSNorm, quant)", 1),
        ([b("Separation: semantics vs implementation vs dispatch")], 0),
        ([("@register_op", {"mono": True}), (" = semantics · ", {}), ("register_impl", {"mono": True}), (" = per-provider kernels (native, vllm_c, aiter, triton_*)", {})], 1),
        ([("Priority lists ", {}), ("--ir-op-priority.<op>=aiter,vllm_c,native", {"mono": True}), (" pick the kernel late", {})], 1),
        ([o("Why it helps fusions"), (": passes match ONE high-level op, not every kernel variant/shape → PR #39242 pattern stays simple", {})], 0),
        ("Eager/compile consistency; maybe_inplace + clone elimination → zero-copy", 0),
        ([o("AMD angle"), (": aiter provider, ROCm default priority lists, register AMD kernels out-of-tree", {})], 0),
    ])

    # S5b: functionalization & maybe_inplace
    add_bullets(prs, "Functionalization & maybe_inplace (why)", [
        ([b("Functionalized graph = no in-place mutation")], 0),
        ("Every op returns NEW tensors → compiler can safely reorder, fuse, pattern-match", 1),
        ("AOTAutograd performs this — it's why all post-grad fusion passes assume a functional graph", 1),
        ([("in-place ", {}), ("rms_norm(out, x, w)", {"mono": True}), ("   →   functional   ", {}), ("y = rms_norm(x, w)", {"mono": True})], 1),
        ([b("The tension:"), (" LLM inference wants IN-PLACE kernels (don't reallocate activations each layer) for memory", {})], 0),
        ([o("maybe_inplace"), (": caller donates inputs — \"overwrite OK, I won't reuse them\"", {})], 0),
        ([("pre-grad ", {}), ("inplace functionalization", {"mono": True}), (": maybe_inplace → functional default, but RECORDS donated inputs", {})], 1),
        ("after lowering to an in-place kernel, clone-elimination drops the safety clones → zero-copy", 1),
        ([o("Benefit:"), (" in-place memory efficiency AND a functional graph the compiler can optimize — identical semantics eager & compiled", {})], 0),
    ])

    # S6: nuances
    add_bullets(prs, "Nuances: control flow & black-box ops", [
        ([b("Full-graph requirement:")], 0),
        ("Forward must capture into one graph dynamic on #tokens", 1),
        ([("Branching on a dynamic shape (", {}), ("if x.size(0) % 128 == 0", {"mono": True}), (") → guards / silent bugs", {})], 1),
        ([b("Two escape hatches")], 0),
        ("Rewrite to avoid the branch", 1),
        ([o("Wrap branchy/opaque logic into a custom op"), (" — Dynamo won't trace inside (a 'black box')", {})], 1),
        ([b("Why register a black-box op")], 0),
        ("Hides un-traceable logic; defines a clean fusion/replacement target; controls split boundary", 1),
        ([("Your fused op ", {}), b("IS"), (" a black-box op (Part III)", {})], 1),
    ])

    # =====================================================================
    # PART III — WORKED EXAMPLE: PR #39242 (MLA Dual RMSNorm Fusion)
    # =====================================================================
    add_section(prs, "Part III — Worked Example: MLA Dual RMSNorm Fusion (PR #39242)")

    # S7: the opportunity
    add_bullets(prs, "The opportunity", [
        ([b("MLA attention runs TWO RMS norms per layer")], 0),
        ([("q_a_layernorm", {"mono": True}), (" (on compressed q) + ", {}), ("kv_a_layernorm", {"mono": True}), (" (on compressed kv)", {})], 1),
        ([("Kimi-K2 = 61 layers → ", {}), b("122 norm kernel launches per forward")], 1),
        ([o("Goal: fuse the pair → ONE AITER fused_qk_rmsnorm HIP kernel"), ("  (2 launches → 1 per layer)", {})], 0),
        ([("Ships as a ", {}), b("compiler pass"), (", not a model edit (philosophy callback)", {})], 0),
        ([("Honest scope: native ", {}), ("rms_norm", {"mono": True}), (" lets Inductor auto-fuse these; this pass targets the", {})], 0),
        ("AITER custom-op case that Inductor cannot fuse on its own", 1),
    ])

    # S8: the pattern in graph terms
    add_code(prs, "The pattern, in graph terms", """
# Unfused FX subgraph (vllm_ir stage)
gemm -> split([q_dim, kv_dim])
   +-- q_c     -> vllm_ir.rms_norm(q_c,  q_w,  eps)            # norm #1
   +-- kv_lora -> split([kv_c_dim, k_pe_dim])
                   +-- kv_c -> vllm_ir.rms_norm(kv_c, kv_w, eps)  # norm #2
                   +-- k_pe                                     # passthrough

# Fused: one connected subgraph -> one op
q_normed, kv_normed = fused_mla_dual_rms_norm(
    q_c, q_w, kv_c, kv_w, eps, eps)        # single AITER HIP kernel
""", note="Matches on vllm_ir.rms_norm — one clean high-level target (Part II payoff)")

    # S9a: files changed overview
    add_bullets(prs, "What the PR changes (8 files)", [
        ([("vllm/_aiter_ops.py", {"mono": True}), (" — register ", {}), ("fused_mla_dual_rms_norm", {"mono": True}), (" custom op (wraps AITER kernel)", {})], 0),
        ([("…/passes/fusion/rocm_aiter_fusion.py", {"mono": True}), (" — pattern + pass via ", {}), ("register_replacement", {"mono": True})], 0),
        ([("vllm/config/compilation.py", {"mono": True}), (" — ", {}), ("fuse_mla_dual_rms_norm", {"mono": True}), (" PassConfig flag + ROCm guard", {})], 0),
        ([("vllm/config/vllm.py", {"mono": True}), (" — optimization-level defaults (auto-enable O1+ on ROCm/AITER)", {})], 0),
        ([("…/passes/pass_manager.py", {"mono": True}), (" — register the pass when enabled", {})], 0),
        ([("tests/compile/passes/test_fuse_mla_dual_rms_norm.py", {"mono": True}), (" — pattern fires, ops replaced, numerics", {})], 0),
        ([("docs/design/fusions.md", {"mono": True}), (" + ", {}), ("optimization_levels.md", {"mono": True}), (" — documentation", {})], 0),
    ])

    # S9b: the fused custom op (black box)
    add_code(prs, "1. The fused custom op (a \"black box\" to Dynamo)", """
# vllm/_aiter_ops.py
def _fused_mla_dual_rms_norm_impl(x1, x1_weight, x2, x2_weight,
                                  x1_epsilon, x2_epsilon):
    import aiter.ops.fused_qk_norm_rope_cache_quant as aiter_ops
    return aiter_ops._fused_qk_rmsnorm(
        q=x1, q_weight=x1_weight, q_eps=x1_epsilon,
        k=x2, k_weight=x2_weight, k_eps=x2_epsilon)

def _fused_mla_dual_rms_norm_fake(x1, x1_weight, x2, x2_weight, e1, e2):
    return (torch.empty_like(x1), torch.empty_like(x2))   # shapes for compile

direct_register_custom_op(
    op_name="fused_mla_dual_rms_norm",
    op_func=_fused_mla_dual_rms_norm_impl,
    mutates_args=[],
    fake_impl=_fused_mla_dual_rms_norm_fake)
""", size=11)

    # S9c: the pattern pass
    add_code(prs, "2. The pattern-matcher pass", """
# vllm/compilation/passes/fusion/rocm_aiter_fusion.py
class MLADualRMSNormPattern(VllmPatternReplacement):
    def pattern(projected, q_weight, kv_weight):
        q_c, kv_lora = projected.split([q_dim, kv_dim], dim=-1)
        kv_c, k_pe   = kv_lora.split([kv_c_dim, k_pe_dim], dim=-1)
        q_normed  = vllm.ir.ops.rms_norm(q_c,  q_weight,  eps)   # MATCH
        kv_normed = vllm.ir.ops.rms_norm(kv_c, kv_weight, eps)
        return q_normed, kv_normed, k_pe

    def replacement(projected, q_weight, kv_weight):
        q_c, kv_lora = projected.split([q_dim, kv_dim], dim=-1)
        kv_c, k_pe   = kv_lora.split([kv_c_dim, k_pe_dim], dim=-1)
        q_normed, kv_normed = torch.ops.vllm.fused_mla_dual_rms_norm(
            q_c, q_weight, kv_c, kv_weight, eps, eps)           # REPLACE
        return q_normed, kv_normed, k_pe

class MLADualRMSNormFusionPass(VllmFusionPatternMatcherPass):
    def __init__(self, config):
        super().__init__(config, "mla_dual_rms_norm_fusion_pass")
        for epsilon in [1e-5, 1e-6]:                  # eps baked into pattern
            self.register(MLADualRMSNormPattern(epsilon))
""", size=10)

    # S9d: wiring & gating
    add_code(prs, "3. Wiring & gating", """
# config/compilation.py — the PassConfig flag (+ ROCm guard)
fuse_mla_dual_rms_norm: bool = None
if self.fuse_mla_dual_rms_norm and not current_platform.is_rocm():
    self.fuse_mla_dual_rms_norm = False        # auto-disable off-ROCm

# config/vllm.py — optimization-level defaults
#   O0: False   |   O1/O2/O3: enable_mla_dual_rms_norm_fusion (ROCm+AITER)

# passes/pass_manager.py — register only when truly available
if (self.pass_config.fuse_mla_dual_rms_norm
        and rocm_aiter_ops.is_enabled()
        and check_aiter_fused_qk_rmsnorm()):
    self.passes += [MLADualRMSNormFusionPass(config)]
""", size=11)

    # S10: verifying
    add_bullets(prs, "Verifying the fusion", [
        ([b("Unit test"), (" — ", {}), ("test_fuse_mla_dual_rms_norm", {"mono": True})], 0),
        ([("before: ", {}), ("vllm_ir.rms_norm", {"mono": True}), ("  →  after: ", {}), ("vllm.fused_mla_dual_rms_norm", {"mono": True})], 1),
        ([("asserts ", {}), ("matched_count == 1", {"mono": True}), (" and numerics close (atol/rtol 1e-2)", {})], 1),
        ([b("Fires on all TP workers:"), (" ", {}), ("MLADualRMSNormFusionPass: fused 1 q/kv norm pair(s)", {"mono": True})], 0),
        ([b("Trace:"), (" two CK RMSNorm kernels → one ", {}), ("aiter::fused_qk_rmsnorm_kernel", {"mono": True})], 0),
        ([b("Accuracy gate:"), (" GSM8K 0.94 ≥ 0.90 (Kimi-K2-Thinking-MXFP4, TP=4, MI355X)", {})], 0),
    ])

    # S10b: perf scorecard
    add_table_slide(
        prs,
        "Measuring — Kimi-K2-Thinking-MXFP4, TP=4, MI355X",
        ["Config", "Concurrency", "Baseline (tok/s/GPU)", "QK Fusion (tok/s/GPU)", "Uplift"],
        [
            ["1k/1k", "16", "521", "533", "1.02x"],
            ["1k/1k", "32", "783", "801", "1.02x"],
            ["8k/1k", "4", "799", "819", "1.03x"],
            ["8k/1k", "32", "2917", "2975", "1.02x"],
        ],
        note="Geomean ~1.02x — small op; the value here is the reusable methodology",
    )

    # =====================================================================
    # PART IV — PRACTICAL: CONFIG & DEBUG
    # =====================================================================
    add_section(prs, "Part IV — Practical: Config & Debug")

    # S11: a real compilation config — code left, explanation right
    add_code_explain(
        prs,
        "Compilation config in practice",
        """
# Enable just the PR's fusion (from PR #39242):
--compilation-config='{"pass_config":
  {"fuse_mla_dual_rms_norm": true}}'

# A fuller real-world config:
--compilation-config '{
  "mode": 3,
  "pass_config": {
    "fuse_allreduce_rms": true,
    "eliminate_noops": true,
    "fuse_rope_kvcache_cat_mla": true
  },
  "custom_ops": ["none"],
  "compile_ranges_endpoints": [64],
  "cudagraph_mode": "FULL_AND_PIECEWISE",
  "use_inductor_graph_partition": true
}'
""",
        [
            ([("mode: 3", {"mono": True}), (" = ", {}), b("VLLM_COMPILE"), (" (vLLM custom backend)", {})], 0),
            ([("pass_config", {"mono": True}), (": toggle individual fusions", {})], 0),
            ("fuse_allreduce_rms · eliminate_noops · fuse_rope_kvcache_cat_mla", 1),
            ([o("fuse_mla_dual_rms_norm"), (" = our PR (ROCm/AITER)", {})], 1),
            ([("custom_ops [\"none\"]", {"mono": True}), (": disable custom ops → Inductor lowers & fuses", {})], 0),
            ("add an op back (e.g. +rms_norm) only when a pattern fusion needs it", 1),
            ([("compile_ranges_endpoints [64]", {"mono": True})], 0),
            ("specialized graph [1,64], general [65, max]", 1),
            ([("cudagraph_mode FULL_AND_PIECEWISE", {"mono": True}), (": full graph (uniform decode) + piecewise (prefill/mixed)", {})], 0),
            ([("use_inductor_graph_partition", {"mono": True}), (": passes see full graph (torch≥2.9)", {})], 0),
        ],
    )

    # S11c: other useful knobs
    add_bullets(prs, "Other useful knobs (quick reference)", [
        ([o("Optimization levels"), (": ", {}), ("-O0", {"mono": True}), (" / ", {}), ("-O1", {"mono": True}), (" / ", {}), ("-O2", {"mono": True}), (" (default) / ", {}), ("-O3", {"mono": True}), (" — presets of the flags below", {})], 0),
        ([("-cc.backend=eager", {"mono": True}), (" — turn off TorchInductor (keep Dynamo)", {})], 0),
        ([("compile_sizes=[1,2,4,8]", {"mono": True}), (" — specialized + autotuned kernels per size", {})], 0),
        ([("cudagraph_capture_sizes=[...]", {"mono": True}), (" — control which sizes get captured", {})], 0),
        ([("--ir-op-priority.rms_norm=aiter,vllm_c,native", {"mono": True}), (" — pick kernel per op", {})], 0),
        ([("dynamic_shapes_config.type", {"mono": True}), (" = ", {}), ("backed", {"mono": True}), (" / ", {}), ("unbacked", {"mono": True}), (" / ", {}), ("backed_size_oblivious", {"mono": True})], 0),
        ([("compile_cache_save_format=unpacked", {"mono": True}), (" — editable Inductor output code", {})], 0),
    ])

    prs.save(OUTPUT)
    print(f"Saved {OUTPUT} with {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    build()
